from pptx import Presentation
from pptx.util import Inches, Pt

def generate_presentation():
    prs = Presentation()

    # --- Slide 1: Title ---
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "🛡️ Brute Force Attack Detector"
    slide.placeholders[1].text = "Machine Learning for Cybersecurity\nCreated by [Your Name]"

    # --- Slide 2: The Problem ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "The Threat: Brute Force"
    content = slide.placeholders[1]
    content.text = ("• Attackers attempt thousands of passwords to gain access.\n"
                    "• Static security rules (e.g., 5-try lockout) can be bypassed.\n"
                    "• Need for intelligent detection of 'behavioral patterns'.")

    # --- Slide 3: How it Works (Logic) ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Machine Learning Logic"
    content = slide.placeholders[1]
    content.text = ("• Algorithm: Random Forest Classifier\n"
                    "• Input Features:\n"
                    "  - Attempt Velocity (Attempts per minute)\n"
                    "  - Failure Ratio (Failures vs. Successes)\n"
                    "  - Target Scope (Unique Usernames targeted)")
    
    

    # --- Slide 4: System Architecture ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "System Architecture"
    content = slide.placeholders[1]
    content.text = ("• Frontend: Streamlit Web Interface\n"
                    "• Backend: Python & Pandas for Data Processing\n"
                    "• Intelligence: Scikit-Learn Predictive Model")

    # --- Slide 5: Conclusion ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Project Impact"
    content = slide.placeholders[1]
    content.text = ("• Automates log analysis for security teams.\n"
                    "• Visualizes attack vectors in real-time.\n"
                    "• Highly scalable for large enterprise logs.")

    # Save the presentation
    filename = "BruteForce_Presentation.pptx"
    prs.save(filename)
    print(f"Success! {filename} has been created.")

if __name__ == "__main__":
    generate_presentation()